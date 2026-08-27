#!/usr/bin/env python3
"""Aniverse 통합 발표 PPT 생성 — 기존 이미지 + 역할분담 + Before/After + 트러블슈팅."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

BASE = Path(__file__).resolve().parent
IMG = BASE / "images"
OUT = BASE / "ppt"
OUT.mkdir(parents=True, exist_ok=True)

BLUE = RGBColor(37, 99, 235)
NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(71, 85, 105)
LIGHT = RGBColor(248, 250, 252)
TEAL = RGBColor(13, 148, 136)
GREEN = RGBColor(22, 163, 74)
AMBER = RGBColor(245, 158, 11)
RED = RGBColor(220, 38, 38)
PURPLE = RGBColor(124, 58, 237)
ORANGE = RGBColor(234, 88, 12)
WHITE = RGBColor(255, 255, 255)


def font(run, size=14, bold=False, color=NAVY):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Malgun Gothic"


def add_header(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(12.4), Inches(0.6))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    font(r, 22, True, WHITE)
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        font(r2, 11, False, RGBColor(191, 219, 254))


def add_footer(slide, text="Aniverse 발표자료 (통합)"):
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.25))
    p = tx.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    font(r, 9, False, SLATE)
    p.alignment = PP_ALIGN.RIGHT


def add_bullets(slide, x, y, w, h, title, bullets, color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = color
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    font(r, 15, True, color)
    for item in bullets:
        bp = tf.add_paragraph()
        br = bp.add_run()
        br.text = item
        font(br, 12, False, NAVY)
        bp.space_after = Pt(3)


def add_table(slide, x, y, w, h, data, header_fill=RGBColor(219, 234, 254)):
    rows, cols = len(data), len(data[0])
    table = slide.shapes.add_table(rows, cols, x, y, w, h).table
    row_h = int(h / rows)
    for r in range(rows):
        table.rows[r].height = row_h
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    font(run, 11 if r else 12, r == 0, NAVY)
                p.alignment = PP_ALIGN.CENTER if r == 0 else PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if r == 0 else (WHITE if r % 2 else LIGHT)


def maybe_image(slide, name, x, y, w=None, h=None):
    path = IMG / name
    if path.exists():
        slide.shapes.add_picture(str(path), x, y, width=w, height=h)
        return True
    return False


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 0 cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(2.5))
    tf = t.text_frame
    r = tf.paragraphs[0].add_run()
    r.text = "Aniverse"
    font(r, 40, True, WHITE)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "온프레미스 4서버 → AWS + Terraform / CI·CD 자동화"
    font(r2, 18, False, RGBColor(147, 197, 253))
    p3 = tf.add_paragraph()
    r3 = p3.add_run()
    r3.text = "https://aniverse.my  ·  팀 역할 분담  ·  HTTPS/WAF  ·  자동 배포"
    font(r3, 14, False, RGBColor(186, 230, 253))

    # 1 story
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "1. 한 줄 스토리", "앱은 그대로, 인프라·배포를 클라우드로")
    add_bullets(
        s,
        Inches(0.5),
        Inches(1.2),
        Inches(6.0),
        Inches(4.8),
        "핵심",
        [
            "Before: Nginx · Django · NFS · MariaDB (수동)",
            "After: ALB/WAF/ASG/RDS/EFS/S3 + Redis + Secrets",
            "자동화: Terraform + GitHub Actions + CodeDeploy",
            "보안: HTTPS(ACM) · WAF · SG · Secrets Manager",
            "운영: SSM · CloudWatch · DB/미디어 복구 자동화",
        ],
        BLUE,
    )
    maybe_image(s, "aniverse-architecture-overview.png", Inches(6.8), Inches(1.2), w=Inches(5.9))
    add_footer(s)

    # 2 before after
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "2. Before → After 매핑", "4서버 역할이 AWS 서비스로 대응")
    add_table(
        s,
        Inches(0.5),
        Inches(1.2),
        Inches(12.3),
        Inches(4.8),
        [
            ["온프레미스", "AWS", "인프라 담당"],
            ["Nginx (DMZ)", "ALB + ACM + WAF + EC2 Nginx", "유민 / 현우"],
            ["Django WAS", "ASG EC2 (Daphne) + Redis", "유민"],
            ["NFS", "EFS + S3 미디어", "윤주"],
            ["MariaDB", "RDS + Secrets Manager", "윤주 / 현우"],
            ["수동 배포 / SSH", "Actions → CodeDeploy / SSM", "현우"],
            ["망·방화벽", "VPC · SG · NAT Instance", "서이"],
        ],
    )
    add_footer(s)

    # 3 roles
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "3. AWS 자동화 팀 역할", "앱 기능 담당과 클라우드 모듈 담당은 분리")
    roles = [
        ("박서이\nNetwork & Security", "network / security / nat\nVPC·SG·NAT t3.micro\n라우팅·outputs", PURPLE),
        ("강유민\nCompute & Traffic", "compute / alb\nALB·ASG·Launch Template\nuser_data·Nginx proxy", ORANGE),
        ("김윤주\nData & Storage", "database / storage\nRDS·EFS·S3\nRemote State", GREEN),
        ("김현우\nDevOps & CI/CD", "environments/dev + cicd\nActions·CodeDeploy·SSM\n조립·HTTPS·WAF", RED),
    ]
    for i, (title, body, color) in enumerate(roles):
        x = Inches(0.35) + i * Inches(3.2)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.25), Inches(3.05), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = color
        card.line.width = Pt(2.25)
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        font(r, 15, True, color)
        p.alignment = PP_ALIGN.CENTER
        for line in body.split("\n"):
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = line
            font(r2, 12, False, NAVY)
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(8)
    add_footer(s)

    # 4 app vs aws
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "4. 앱 담당 ↔ 클라우드 담당", "같은 팀, 다른 축의 역할")
    add_table(
        s,
        Inches(1.2),
        Inches(1.5),
        Inches(10.8),
        Inches(4.2),
        [
            ["이름", "앱 (온프레미스)", "AWS 자동화"],
            ["김현우", "auth / anime / 챗봇 / 서버", "DevOps & CI/CD"],
            ["박서이", "deal / 채팅", "Network & Security"],
            ["강유민", "works 창작", "Compute & Traffic"],
            ["김윤주", "community", "Data & Storage"],
        ],
        header_fill=RGBColor(237, 233, 254),
    )
    add_footer(s)

    # 5 architecture image
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "5. 전체 아키텍처", "Public ALB · Private App/DB · S3/EFS · Secrets · Redis")
    maybe_image(s, "aniverse-architecture-overview.png", Inches(0.55), Inches(1.05), w=Inches(12.2))
    add_footer(s)

    # 6 modules
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "6. Terraform 모듈 구성", "역할별 모듈 → environments/dev 에서 조립")
    maybe_image(s, "aniverse-terraform-modules.png", Inches(0.4), Inches(1.05), w=Inches(7.0))
    add_table(
        s,
        Inches(7.6),
        Inches(1.15),
        Inches(5.2),
        Inches(5.0),
        [
            ["담당", "모듈"],
            ["서이", "network / security / nat"],
            ["윤주", "database / storage"],
            ["유민", "alb / compute"],
            ["현우", "dev 조립 · cicd · acm · waf"],
            ["공통", "endpoints · secrets · redis"],
        ],
        header_fill=RGBColor(204, 251, 241),
    )
    add_footer(s)

    # 7 cicd
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "7. CI/CD 파이프라인", "infra 저장소와 app 저장소 분리")
    maybe_image(s, "aniverse-cicd-pipeline.png", Inches(0.4), Inches(1.05), w=Inches(7.0))
    add_table(
        s,
        Inches(7.6),
        Inches(1.15),
        Inches(5.2),
        Inches(3.2),
        [
            ["구분", "내용"],
            ["Infra CD", "terraform apply"],
            ["App CD", "S3 zip → CodeDeploy"],
            ["운영", "SSM · CloudWatch"],
            ["도메인", "Route53 + ACM HTTPS"],
        ],
    )
    add_bullets(
        s,
        Inches(7.6),
        Inches(4.55),
        Inches(5.2),
        Inches(1.8),
        "한 줄",
        ["push → 인프라/앱 자동 반영", "aniverse.my 로 서비스"],
        GREEN,
    )
    add_footer(s)

    # 8 troubleshooting
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "8. 트러블슈팅 · 개선", "실제 장애를 파이프라인/코드로 흡수")
    maybe_image(s, "aniverse-problems-solved.png", Inches(0.35), Inches(1.05), w=Inches(6.6))
    add_table(
        s,
        Inches(7.15),
        Inches(1.1),
        Inches(5.7),
        Inches(5.3),
        [
            ["문제", "해결"],
            ["CD 미동작 / TF버전", "yml + TF 1.10.5"],
            ["ALB 혼재", "alb 모듈 + moved"],
            ["챗봇/WS", "CSRF·Gemini·Daphne"],
            ["이미지 403", "S3 + bucket env"],
            ["CodeDeploy/health", "agent·/health/"],
            ["DB restore env", "Secrets fallback"],
            ["static 403", "home 권한 수정"],
            ["HTTPS", "ACM + ALB + WAF"],
        ],
        header_fill=RGBColor(254, 242, 242),
    )
    add_footer(s)

    # 9 closing
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "9. 마무리", "30초 멘트")
    add_bullets(
        s,
        Inches(0.7),
        Inches(1.3),
        Inches(11.8),
        Inches(4.8),
        "발표 멘트",
        [
            "“앱은 기능 담당으로 만들었고, 클라우드 이전에서는 역할을 다시 나눴습니다.”",
            "“서이: 네트워크·보안 / 유민: ALB·ASG / 윤주: RDS·EFS·S3”",
            "“현우: 모듈 조립 + Terraform·GitHub Actions 자동화”",
            "“HTTPS·WAF·Secrets까지 넣었고, aniverse.my 로 서비스 중입니다.”",
        ],
        ORANGE,
    )
    add_footer(s, "Aniverse 발표자료 (통합) · End")

    out = OUT / "Aniverse_발표_통합.pptx"
    prs.save(out)
    # keep legacy names pointing to same story via copies
    prs.save(OUT / "Aniverse_발표.pptx")
    print("Wrote", out)


if __name__ == "__main__":
    build()
