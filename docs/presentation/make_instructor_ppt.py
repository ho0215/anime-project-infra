#!/usr/bin/env python3
"""강사용 5~7분 발표 PPT — AWS/GitHub 공식 아이콘 구성도 포함."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# 한국어 Windows 기본 글꼴. latin + ea(동아시아) 둘 다 지정해야 한글이 깨지지 않음.
FONT_NAME = "맑은 고딕"

BASE = Path(__file__).resolve().parent
IMG = BASE / "images" / "instructor"
OUT = BASE / "ppt"
OUT.mkdir(parents=True, exist_ok=True)

NAVY = RGBColor(15, 23, 42)
WHITE = RGBColor(255, 255, 255)
SLATE = RGBColor(71, 85, 105)
BLUE = RGBColor(37, 99, 235)
LIGHT = RGBColor(248, 250, 252)
PURPLE = RGBColor(124, 58, 237)
ORANGE = RGBColor(234, 88, 12)
GREEN = RGBColor(22, 163, 74)
RED = RGBColor(220, 38, 38)
TEAL = RGBColor(13, 148, 136)
SOFT = RGBColor(239, 246, 255)


def font(run, size=16, bold=False, color=NAVY):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT_NAME
    # python-pptx는 font.name이 latin만 설정함 → 한글은 a:ea typeface 필요
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(qn(f"a:{tag}"))
        if el is None:
            el = rPr.makeelement(qn(f"a:{tag}"), {})
            rPr.append(el)
        el.set("typeface", FONT_NAME)


def pad_tf(tf, left=0.12, top=0.08, right=0.12, bottom=0.08):
    tf.word_wrap = True
    tf.margin_left = Inches(left)
    tf.margin_right = Inches(right)
    tf.margin_top = Inches(top)
    tf.margin_bottom = Inches(bottom)


def set_text(tf, text, size=16, bold=False, color=NAVY, align=None):
    tf.clear()
    pad_tf(tf)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    font(r, size, bold, color)
    if align is not None:
        p.alignment = align


def add_para(tf, text, size=14, bold=False, color=NAVY, space_before=6):
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = text
    font(r, size, bold, color)
    p.space_before = Pt(space_before)
    return p


def rect(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1.5)
    else:
        sh.line.fill.background()
    pad_tf(sh.text_frame)
    return sh


def header(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.12), Inches(12.4), Inches(0.65))
    tf = box.text_frame
    set_text(tf, title, 22, True, WHITE)
    if subtitle:
        add_para(tf, subtitle, 12, False, RGBColor(147, 197, 253), 2)


TOTAL = 18


def footer(slide, n, total=TOTAL):
    tx = slide.shapes.add_textbox(Inches(0.45), Inches(7.05), Inches(12.4), Inches(0.3))
    p = tx.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = f"Aniverse · 온프레미스 3-tier → AWS   |   {n}/{total}"
    font(r, 10, False, SLATE)
    p.alignment = PP_ALIGN.RIGHT


def put_img(slide, name, x, y, w=None, h=None, max_bottom=Inches(6.95)):
    """이미지를 넣되 footer 위로 잘리지 않게 max_bottom 안에 맞춤."""
    path = IMG / name
    if not path.exists():
        return False
    from PIL import Image

    iw, ih = Image.open(path).size
    aspect = ih / float(iw)
    x_in = x.inches
    y_in = y.inches
    max_h_in = max_bottom.inches - y_in
    if max_h_in <= 0.5:
        return False

    if w is not None and h is not None:
        w_in, h_in = w.inches, h.inches
    elif w is not None:
        w_in = w.inches
        h_in = w_in * aspect
    elif h is not None:
        h_in = h.inches
        w_in = h_in / aspect
    else:
        w_in = 12.4
        h_in = w_in * aspect

    if h_in > max_h_in:
        h_in = max_h_in
        w_in = h_in / aspect
        # 가로로 줄어든 만큼 가운데 정렬
        box_w = w.inches if w is not None else 12.4
        if w_in < box_w:
            x_in = x_in + (box_w - w_in) / 2.0

    slide.shapes.add_picture(
        str(path), Inches(x_in), Inches(y_in), width=Inches(w_in), height=Inches(h_in)
    )
    return True


def cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(12), Inches(3.8))
    tf = t.text_frame
    set_text(tf, "Aniverse", 40, True, WHITE)
    add_para(tf, "온프레미스 3-Tier 서비스 → AWS 이전", 24, False, RGBColor(191, 219, 254), 14)
    add_para(tf, "Terraform · Actions · CI/CD · HA · HTTPS · WAF", 16, False, RGBColor(147, 197, 253), 16)
    add_para(tf, "https://aniverse.my", 15, False, RGBColor(125, 211, 252), 20)


def agenda(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "목차")
    items = [
        "01  AWS 이전하면서 맡은 역할",
        "02  온프레미스 운영에서 어려웠던 점",
        "03  Terraform으로 구성한 인프라",
        "04  GitHub Actions 자동화",
        "05  CI/CD · 고가용성 · 모니터링 · 보안",
        "06  HTTPS · WAF",
        "07  사용한 AWS · GitHub 한 장 정리",
    ]
    for i, title in enumerate(items):
        y = Inches(1.1) + i * Inches(0.78)
        card = rect(s, Inches(0.7), y, Inches(11.9), Inches(0.68), LIGHT, BLUE)
        set_text(card.text_frame, title, 17, True, NAVY)
    footer(s, 2)


def roles(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "1. AWS 이전 — 맡은 역할", "AWS 공식 아이콘으로 담당 서비스만 빠르게 보기")
    put_img(s, "09_team_roles.png", Inches(0.25), Inches(0.95), w=Inches(12.8))
    footer(s, 3)


def pain_detail(prs):
    """02 어려움 — 더 자세히, 이미지(온프렘 구조) 함께."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "2. 온프레미스에서 어려웠던 점 (1)", "3-Tier: Nginx · Django · NFS/MariaDB 를 서버별로 직접 운영")
    put_img(s, "01_onprem_3tier.png", Inches(0.4), Inches(1.1), w=Inches(12.5))
    footer(s, 4)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "2. 온프레미스에서 어려웠던 점 (2)", "수동 운영이 만든 반복 문제")
    pains = [
        (
            "① 배포가 전부 수작업",
            [
                "서버 4대에 SSH로 접속해 pull · restart",
                "순서가 조금만 달라져도 결과가 달라짐",
                "누가 언제 무엇을 배포했는지 기록이 남기 어려움",
            ],
        ),
        (
            "② 환경 불일치",
            [
                "로컬에서는 되는데 서버에서만 실패하는 경우 빈번",
                "패키지·환경변수·권한 설정이 사람 기억에 의존",
                "원인 파악에 시간을 대부분 소모",
            ],
        ),
        (
            "③ 장애·확장 대응이 느림",
            [
                "트래픽 증가 시 서버를 직접 추가·설정",
                "한 대 장애가 곧 서비스 전체 영향으로 이어짐",
                "복구 절차가 문서화되어 있어도 재현이 어려움",
            ],
        ),
        (
            "④ 보안·비밀값 관리",
            [
                ".env, API Key가 서버 파일에 흩어져 존재",
                "HTTPS·방화벽 규칙을 서버마다 따로 맞춤",
                "권한/키 교체 시 누락 위험이 큼",
            ],
        ),
    ]
    for i, (title, bullets) in enumerate(pains):
        x = Inches(0.35) + (i % 2) * Inches(6.45)
        y = Inches(1.15) + (i // 2) * Inches(2.8)
        card = rect(s, x, y, Inches(6.2), Inches(2.6), LIGHT, RED)
        tf = card.text_frame
        set_text(tf, title, 15, True, RED)
        for b in bullets:
            add_para(tf, "• " + b, 12, False, NAVY, 5)
    footer(s, 5)


def terraform_slides(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "3. Terraform 인프라 구성 (1)", "온프레미스 3-Tier를 AWS 서비스로 대응")
    put_img(s, "02_aws_overview.png", Inches(0.35), Inches(1.05), w=Inches(12.6))
    footer(s, 6)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "3. Terraform 인프라 구성 (2)", "모듈을 나눠 만들고 environments/dev 에서 조립")
    put_img(s, "03_terraform_modules.png", Inches(0.25), Inches(1.0), w=Inches(12.8))
    footer(s, 7)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "3. Terraform 인프라 구성 (3)", "주요 리소스와 연결 관계")
    left = rect(s, Inches(0.35), Inches(1.15), Inches(6.25), Inches(5.55), SOFT, BLUE)
    tf = left.text_frame
    set_text(tf, "계층별 구성", 16, True, BLUE)
    for prefix, line in [
        ("• ", "Network: VPC · Public/Private Subnet · IGW · NAT"),
        ("• ", "Security: ALB / App / NAT / DB / EFS"),
        ("    ", "Redis / endpoints 용 Security Group"),
        ("• ", "Traffic: ALB · Target Group · HTTPS Listener"),
        ("• ", "Compute: Launch Template · ASG · IAM Role"),
        ("• ", "Data: RDS MariaDB · EFS · S3"),
        ("• ", "Ops: Secrets Manager · WAF · Redis · CodeDeploy"),
    ]:
        add_para(tf, prefix + line, 12, False, NAVY, 7)

    right = rect(s, Inches(6.85), Inches(1.15), Inches(6.1), Inches(5.55), LIGHT, TEAL)
    tf = right.text_frame
    set_text(tf, "조립 방식", 16, True, TEAL)
    for line in [
        "1) S3 + DynamoDB로 Remote State 환경 선구축",
        "2) 팀원이 모듈 PR로 기능 단위 작성",
        "3) outputs.tf 로 vpc_id, sg_id 등 연결",
        "4) environments/dev/main.tf 에서 module 호출",
        "5) terraform plan / apply 로 한 번에 반영",
        "",
        "결과",
        "• 인프라가 코드로 리뷰·재현 가능",
        "• 서버 ‘기억 의존’ 감소",
        "• 팀 전체가 충돌 없이 동시 작업 가능 (State Lock)",
    ]:
        add_para(tf, line, 12, False, NAVY, 5)
    footer(s, 8)

    # 데이터베이스·스토리지 구성 요소 (AWS 아이콘)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "3. 데이터베이스·스토리지 구성", "김윤주 Data & Storage — RDS · EFS · S3 (AWS 공식 아이콘)")
    # reuse storage roles visual for this slide as overview cards appear in image
    put_img(s, "08_storage_roles.png", Inches(0.25), Inches(0.95), w=Inches(12.8))
    footer(s, 9)

    # S3 / EFS / RDS — 저장 역할 (같은 이미지 확대 설명용 텍스트 보조 슬라이드)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "3. S3 · EFS · RDS — 저장 역할 정리", "RDS=구조화 데이터 · EFS=공유 파일 · S3=객체/배포")
    left = rect(s, Inches(0.35), Inches(1.15), Inches(4.0), Inches(5.55), LIGHT, GREEN)
    tf = left.text_frame
    set_text(tf, "RDS (MariaDB)", 16, True, GREEN)
    for line in [
        "• 회원 · 로그인 계정",
        "• 커뮤니티 글 · 댓글",
        "• 거래 · 채팅 메타",
        "• 창작물 제목·본문",
        "• 관계·트랜잭션 필요 값",
        "",
        "modules/database",
        "Private subnet · Secrets",
    ]:
        add_para(tf, line, 13, False, NAVY, 6)

    mid = rect(s, Inches(4.65), Inches(1.15), Inches(4.0), Inches(5.55), LIGHT, TEAL)
    tf = mid.text_frame
    set_text(tf, "EFS", 16, True, TEAL)
    for line in [
        "• EC2 media/ NFS 마운트",
        "• 인스턴스 교체돼도 유지",
        "• 여러 EC2가 같은 media",
        "• 온프렘 NFS 대체",
        "• Nginx /media 서빙",
        "",
        "modules/storage",
        "ASG 공유 파일 공간",
    ]:
        add_para(tf, line, 13, False, NAVY, 6)

    right = rect(s, Inches(8.95), Inches(1.15), Inches(4.0), Inches(5.55), LIGHT, ORANGE)
    tf = right.text_frame
    set_text(tf, "S3", 16, True, ORANGE)
    for line in [
        "• 업로드 이미지 원본",
        "• community / goods / works",
        "• 공개 URL 제공",
        "• 배포 zip 패키지",
        "• tfstate 버킷 (별도)",
        "",
        "CORS · Lifecycle",
        "DynamoDB State Lock",
    ]:
        add_para(tf, line, 13, False, NAVY, 6)
    footer(s, 10)


def actions_slides(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "4. GitHub Actions (1)", "GitHub · AWS 공식 아이콘으로 본 CI/CD Pipeline")
    put_img(s, "04_github_actions.png", Inches(0.25), Inches(0.95), w=Inches(12.8))
    footer(s, 11)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "4. GitHub Actions (2)", "실제로 돌아가는 파이프라인")
    left = rect(s, Inches(0.35), Inches(1.15), Inches(6.25), Inches(5.55), SOFT, BLUE)
    tf = left.text_frame
    set_text(tf, "anime-project-infra", 15, True, BLUE)
    for line in [
        "trigger: main push / workflow_dispatch",
        "steps: checkout → setup TF → fmt",
        "        validate → plan → apply",
        "결과: VPC/ALB/ASG/RDS/S3/WAF 갱신",
        "Secrets: AWS 자격증명, TF_VAR_*",
    ]:
        add_para(tf, "• " + line, 12, False, NAVY, 8)

    right = rect(s, Inches(6.85), Inches(1.15), Inches(6.1), Inches(5.55), RGBColor(255, 247, 237), ORANGE)
    tf = right.text_frame
    set_text(tf, "anime-project", 15, True, ORANGE)
    for line in [
        "trigger: main push / workflow_dispatch",
        "steps: media → S3 sync",
        "        앱 zip 패키징 → deploy S3",
        "        CodeDeploy → ASG 배포",
        "hooks: install → migrate → Daphne 기동",
    ]:
        add_para(tf, "• " + line, 12, False, NAVY, 8)
    footer(s, 12)


def ops_slides(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "5. CI/CD · HA · 모니터링 · 보안", "AWS·GitHub 공식 아이콘으로 한 장 요약")
    put_img(s, "05_ops_security.png", Inches(0.25), Inches(0.95), w=Inches(12.8))
    footer(s, 13)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "5. 적용 내용 상세", "서비스가 실제로 어떻게 버티는지")
    items = [
        ("CI/CD", BLUE, ["Actions + CodeDeploy로 배포 경로 고정", "ALB /health/ 로 배포 성공 판정", "실패 시 Actions·배포 로그로 추적"]),
        ("고가용성", ORANGE, ["ALB로 트래픽 분산", "ASG로 인스턴스 교체·확장", "App/DB 서브넷 분리 배치"]),
        ("모니터링", TEAL, ["CloudWatch 알람", "Target Group health 확인", "SSM으로 서버 점검"]),
        ("보안", RED, ["SG 최소 포트 개방", "DB subnet NACL 설정", "Secrets Manager로 키 주입", "HTTPS(암호화) + WAF(공격 차단)"]),
    ]
    for i, (title, color, bullets) in enumerate(items):
        x = Inches(0.35) + (i % 2) * Inches(6.45)
        y = Inches(1.15) + (i // 2) * Inches(2.85)
        card = rect(s, x, y, Inches(6.2), Inches(2.7), LIGHT, color)
        tf = card.text_frame
        set_text(tf, title, 16, True, color)
        for b in bullets:
            add_para(tf, "• " + b, 12, False, NAVY, 5)
    footer(s, 14)


def https_waf_slides(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "6. HTTPS · WAF (1)", "Route 53 · WAF · ACM · ALB · EC2 (AWS 공식 아이콘)")
    put_img(s, "06_https_waf.png", Inches(0.25), Inches(0.95), w=Inches(12.8))
    footer(s, 15)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "6. HTTPS · WAF (2)", "용어 말고 ‘무엇을 했는지’만")

    left = rect(s, Inches(0.35), Inches(1.15), Inches(6.2), Inches(5.55), SOFT, TEAL)
    tf = left.text_frame
    set_text(tf, "HTTPS — 통신 암호화", 17, True, TEAL)
    add_para(tf, "쉽게: 사이트와 사용자 사이 내용을 남이 못 보게 잠근다", 12, True, SLATE, 8)
    for line in [
        "• ACM = 인증서(자물쇠) 발급 서비스",
        "• aniverse.my 용 인증서를 ALB에 붙임",
        "• https(443)로 접속 받기",
        "• http로 들어오면 https로 자동 이동",
        "• 결과: https://aniverse.my + 자물쇠",
    ]:
        add_para(tf, line, 14, False, NAVY, 10)

    right = rect(s, Inches(6.8), Inches(1.15), Inches(6.15), Inches(5.55), RGBColor(254, 242, 242), RED)
    tf = right.text_frame
    set_text(tf, "WAF — 앞단 문지기", 17, True, RED)
    add_para(tf, "쉽게: 서버 들어가기 전에 나쁜 요청을 막는다", 12, True, SLATE, 8)
    for line in [
        "• WAF = 웹 방화벽 (ALB 앞)",
        "• SQL 삽입·이상한 입력 차단",
        "• 한 IP가 너무 자주 치면 차단",
        "• 통과한 요청만 EC2로 전달",
        "• 결과: 공격성 트래픽을 입구에서 감소",
    ]:
        add_para(tf, line, 14, False, NAVY, 10)
    footer(s, 16)


def stack_summary_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "7. 사용한 AWS · GitHub", "인프라 서비스와 자동화 도구를 한 장으로")
    put_img(s, "07_aws_github_stack.png", Inches(0.35), Inches(1.0), w=Inches(12.6))
    footer(s, 17)


def closing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4))
    tf = t.text_frame
    set_text(tf, "정리", 20, True, RGBColor(147, 197, 253))
    add_para(tf, "3-Tier 수동 서버 → AWS 모듈형 인프라 + 자동 배포", 22, True, WHITE, 14)
    add_para(tf, "서이(망) · 유민(트래픽) · 윤주(데이터) · 현우(자동화·HTTPS·WAF)", 16, False, RGBColor(191, 219, 254), 16)
    add_para(tf, "https://aniverse.my", 16, False, RGBColor(125, 211, 252), 14)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    cover(prs)
    agenda(prs)
    roles(prs)
    pain_detail(prs)
    terraform_slides(prs)
    actions_slides(prs)
    ops_slides(prs)
    https_waf_slides(prs)
    stack_summary_slide(prs)
    closing(prs)
    out = OUT / "Aniverse_발표_강사용.pptx"
    prs.save(out)
    print("Wrote", out)


if __name__ == "__main__":
    build()
