#!/usr/bin/env python3
"""Generate a simple Aniverse presentation PPTX (no draw.io needed)."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE

OUT = Path(__file__).resolve().parent / "ppt"
OUT.mkdir(parents=True, exist_ok=True)

NAVY = RGBColor(15, 23, 42)
WHITE = RGBColor(255, 255, 255)
SLATE = RGBColor(71, 85, 105)
BLUE = RGBColor(37, 99, 235)
PURPLE = RGBColor(124, 58, 237)
ORANGE = RGBColor(234, 88, 12)
GREEN = RGBColor(22, 163, 74)
RED = RGBColor(220, 38, 38)
LIGHT = RGBColor(248, 250, 252)


def set_run(p, text, size=18, bold=False, color=NAVY):
    p.clear()
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Malgun Gothic"


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.5), Inches(2))
    tf = box.text_frame
    set_run(tf.paragraphs[0], title, 36, True, WHITE)
    p2 = tf.add_paragraph()
    set_run(p2, subtitle, 16, False, RGBColor(147, 197, 253))
    return slide


def add_bullets_slide(prs, title, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLUE
    bg.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
    set_run(t.text_frame.paragraphs[0], title, 24, True, WHITE)

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(11.5), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        set_run(p, line, 18, False, NAVY)
        p.space_after = Pt(10)
    return slide


def add_role_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLUE
    bg.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
    set_run(t.text_frame.paragraphs[0], "AWS 자동화 팀 역할 분담", 24, True, WHITE)

    roles = [
        ("박서이\nNetwork & Security", "network / security / nat\nVPC·SG·NAT Instance\n망 분리·라우팅·outputs", PURPLE),
        ("강유민\nCompute & Traffic", "compute / alb\nALB·ASG·Launch Template\nuser_data·Nginx proxy", ORANGE),
        ("김윤주\nData & Storage", "database / storage\nRDS·EFS·S3\nRemote State 초기화", GREEN),
        ("김현우\nDevOps & CI/CD", "environments/dev\nActions·CodeDeploy·SSM\n모듈 조립·알람", RED),
    ]
    x0 = Inches(0.35)
    w = Inches(3.0)
    gap = Inches(0.2)
    for i, (title, body, color) in enumerate(roles):
        x = x0 + i * (w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.3), w, Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = color
        card.line.width = Pt(2.5)
        tf = card.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        set_run(p, title, 16, True, color)
        p.alignment = PP_ALIGN.CENTER
        for line in body.split("\n"):
            p2 = tf.add_paragraph()
            set_run(p2, line, 13, False, NAVY)
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(8)
    return slide


def add_table_slide(prs, title, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLUE
    bg.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
    set_run(t.text_frame.paragraphs[0], title, 24, True, WHITE)

    cols = len(rows[0])
    table_shape = slide.shapes.add_table(len(rows), cols, Inches(0.5), Inches(1.3), Inches(12.2), Inches(0.55 * len(rows)))
    table = table_shape.table
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(13)
                    run.font.name = "Malgun Gothic"
                    run.font.bold = r == 0
                    run.font.color.rgb = WHITE if r == 0 else NAVY
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, "Aniverse", "온프레미스 4서버 → AWS + Terraform / CI·CD 자동화")
    add_bullets_slide(prs, "한 줄 스토리", [
        "• Before: Nginx · Django · NFS · MariaDB (서버 4대, 수동 운영)",
        "• After: ALB/WAF/ASG/RDS/EFS/S3 + Secrets + Redis",
        "• 자동화: Terraform(IaC) + GitHub Actions + CodeDeploy",
        "• 서비스: https://aniverse.my (HTTPS)",
    ])
    add_table_slide(prs, "Before → After 매핑", [
        ["온프레미스", "AWS"],
        ["Nginx 서버", "ALB + ACM + WAF + EC2 Nginx"],
        ["Django 서버", "ASG EC2 (Daphne) + Redis"],
        ["NFS 서버", "EFS + S3 미디어"],
        ["MariaDB 서버", "RDS + Secrets Manager"],
        ["수동 배포", "Actions → CodeDeploy / SSM"],
    ])
    add_role_slide(prs)
    add_table_slide(prs, "앱 담당 vs 클라우드 담당", [
        ["이름", "앱(온프레미스)", "AWS 자동화"],
        ["김현우", "auth / anime / 챗봇 / 서버", "DevOps & CI/CD"],
        ["박서이", "deal / 채팅", "Network & Security"],
        ["강유민", "works 창작", "Compute & Traffic"],
        ["김윤주", "community", "Data & Storage"],
    ])
    add_bullets_slide(prs, "발표 멘트 (30초)", [
        "“앱은 그대로 두고, 인프라는 네 명이 나눠 맡았습니다.”",
        "“서이: 네트워크·보안 / 유민: ALB·ASG / 윤주: RDS·EFS·S3”",
        "“현우: 모듈 조립 + Terraform·GitHub Actions 자동화”",
        "“결과: aniverse.my 에서 HTTPS + 자동 배포로 서비스 중”",
    ])

    out = OUT / "Aniverse_발표.pptx"
    prs.save(out)
    print(f"Wrote {out}")


if __name__ == "__main__":
    main()
